"""スライド内容の構造データ(list[dict])から実際の.pptxファイルを組み立てる。

デザインの作り込みは行わず、タイトル+箇条書き+スピーカーノートという
シンプルな構成にしている。見た目の調整は人間がPowerPoint/Resolve側の
作業で行う前提。

slide_dataに "image_path" が設定されている場合は、箇条書きの右側に
その画像(Geminiで生成した挿絵など)を配置する。
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

TITLE_LAYOUT_INDEX = 0  # タイトルスライド
CONTENT_LAYOUT_INDEX = 1  # タイトル+コンテンツ


def build_pptx(title: str, slides: list[dict], output_path: str | Path) -> Path:
    """slidesの各要素 {"title", "bullets", "notes", "image_path"(任意)} からpptxを生成する。"""
    prs = Presentation()

    # 表紙
    title_slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT_INDEX])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "VOICEVOX解説動画 スライド"

    # 各コンテンツスライド
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT_INDEX])
        slide.shapes.title.text = slide_data.get("title", "")

        body = slide.placeholders[1]
        image_path = slide_data.get("image_path")

        # 画像がある場合は箇条書きの幅を狭めて右側に画像を置くスペースを空ける
        if image_path:
            body.width = Emu(int(prs.slide_width * 0.55))

        text_frame = body.text_frame
        text_frame.clear()

        bullets = slide_data.get("bullets", [])
        for i, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(22)

        if image_path and Path(image_path).exists():
            image_left = Emu(int(prs.slide_width * 0.60))
            image_top = Inches(1.8)
            image_width = Emu(int(prs.slide_width * 0.35))
            slide.shapes.add_picture(
                str(image_path), image_left, image_top, width=image_width
            )

        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
